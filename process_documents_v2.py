import os
import json
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
import PyPDF2
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from collections import OrderedDict

# Lataa ympäristömuuttujat
load_dotenv()

class DocumentProcessor:
    def __init__(self):
        self.client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
        self.docs_dir = Path('docs')
        self.template_path = Path('master-document-template.json')
        self.output_json = Path('consolidated_due_diligence.json')
        self.output_docx = Path('consolidated_due_diligence.docx')

    def read_pdf(self, file_path):
        """Lukee PDF-tiedoston ja palauttaa tekstin"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"Virhe PDF:n lukemisessa {file_path}: {e}")
            return ""

    def read_excel(self, file_path):
        """Lukee Excel-tiedoston ja palauttaa tekstimuotoisen esityksen"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n=== Taulukko: {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            return text
        except Exception as e:
            print(f"Virhe Excel:in lukemisessa {file_path}: {e}")
            return ""

    def read_pptx(self, file_path):
        """Lukee PowerPoint-tiedoston ja palauttaa tekstin"""
        try:
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== Dia {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Virhe PowerPoint:in lukemisessa {file_path}: {e}")
            return ""

    def read_document(self, file_path):
        """Lukee dokumentin sen tyypin mukaan"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension == '.pdf':
            return self.read_pdf(file_path)
        elif extension in ['.xlsx', '.xls']:
            return self.read_excel(file_path)
        elif extension in ['.pptx', '.ppt']:
            return self.read_pptx(file_path)
        else:
            print(f"Tuntematon tiedostotyyppi: {extension}")
            return ""

    def load_template(self):
        """Lataa master template"""
        with open(self.template_path, 'r', encoding='utf-8') as f:
            return json.load(f, object_pairs_hook=OrderedDict)

    def get_section_info(self, section_key, section_data):
        """Palauttaa osion tiedot tiiviissä muodossa"""
        if not isinstance(section_data, dict):
            return ""

        info = f"\n## {section_key} ##\n"
        for sub_key, sub_data in section_data.items():
            if isinstance(sub_data, dict):
                instruction = sub_data.get('instruction', '')
                update_rule = sub_data.get('update_rule', '')
                if instruction:
                    info += f"  - {sub_key}:\n"
                    info += f"    Update rule: {update_rule}\n"
                    info += f"    Instruction: {instruction}\n"
        return info

    def extract_section_from_document(self, document_text, document_name, section_key, section_template):
        """Poimii yhden osion tiedot dokumentista"""

        # Jos osio on locked, ohita se
        if isinstance(section_template, dict):
            # Tarkista onko päätason locked
            if section_template.get('update_rule') == 'locked':
                return None

            # Tarkista onko kaikki alaosiot locked
            all_locked = True
            for sub_key, sub_data in section_template.items():
                if isinstance(sub_data, dict):
                    if sub_data.get('update_rule') != 'locked':
                        all_locked = False
                        break

            if all_locked:
                return None

        section_info = self.get_section_info(section_key, section_template)

        if not section_info.strip():
            return None

        prompt = f"""Olet tarkka dokumenttianalyytikko. Analysoi dokumentti ja poimii VAIN seuraavan osion tiedot:

{section_info}

KRIITTISET SÄÄNNÖT:
- ÄLÄ KEKSI mitään tietoa
- Palauta VAIN ne alakohdat, joihin dokumentissa on eksplisiittistä tietoa
- Jos et löydä tietoa johonkin alakohtaan, älä sisällytä sitä vastaukseen
- Noudata tarkkaan kunkin kentän "instruction"-ohjeita

Dokumentin nimi: {document_name}

Dokumentin sisältö:
{document_text}

HUOMAA:
- Lue KOKO dokumentti huolellisesti
- Etsi kaikki mahdolliset tiedot jotka liittyvät tähän osioon
- Älä jätä mitään pois
- Jos tieto vaikuttaa liittyvän osioon, sisällytä se

Palauta JSON-objekti, jossa avaimet ovat alakohtien nimiä ja arvot ovat poimittuja tietoja.

Esimerkki paluuarvosta:
{{
  "2_1_company_description": {{
    "value": "Yritys tekee X ja Y",
    "confidence": "high"
  }},
  "2_2_problem_statement": {{
    "value": "Asiakkaat kärsivät ongelmasta Z",
    "confidence": "medium"
  }}
}}

Palauta VAIN JSON, ei muuta tekstiä."""

        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "Olet tarkka dokumenttianalyytikko. Palauta VAIN JSON, ei muuta tekstiä. Jos et löydä tietoa, palauta tyhjä objekti {}."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.05,
                max_tokens=4096
            )

            result = response.choices[0].message.content.strip()

            # Etsi JSON-lohko
            if "```json" in result:
                result = result.split("```json")[1].split("```")[0].strip()
            elif "```" in result:
                result = result.split("```")[1].split("```")[0].strip()

            try:
                extracted = json.loads(result)
                return extracted if extracted else None
            except json.JSONDecodeError as e:
                print(f"  JSON-parsinta epaonnistui osiossa {section_key}: {e}")
                return None

        except Exception as e:
            print(f"  Virhe API-kutsussa osiossa {section_key}: {e}")
            return None

    def build_consolidated_document(self, all_sections_data, template):
        """Rakentaa konsolidoidun dokumentin säilyttäen järjestyksen"""
        result = OrderedDict()

        # Käy läpi template järjestyksessä
        for section_key in template.keys():
            if section_key == 'document_metadata':
                # Metadata käsitellään erikseen lopussa
                continue

            section_content = OrderedDict()

            # Kerää kaikki dokumenttien tiedot tälle osiolle
            for doc_name, sections in all_sections_data.items():
                if section_key not in sections or sections[section_key] is None:
                    continue

                section_data = sections[section_key]
                template_section = template[section_key]

                # Käsittele jokainen alakohta
                for sub_key, sub_value in section_data.items():
                    if not isinstance(sub_value, dict):
                        continue

                    value = sub_value.get('value', '')
                    confidence = sub_value.get('confidence', 'medium')

                    if not value:
                        continue

                    # Hae update_rule templatesta
                    update_rule = 'append'
                    if isinstance(template_section, dict) and sub_key in template_section:
                        sub_template = template_section[sub_key]
                        if isinstance(sub_template, dict):
                            update_rule = sub_template.get('update_rule', 'append')

                    # Älä käsittele locked-kenttiä
                    if update_rule == 'locked':
                        continue

                    # Käsittele update_rule
                    if update_rule == 'overwrite':
                        section_content[sub_key] = value
                    elif update_rule == 'append':
                        if sub_key not in section_content:
                            section_content[sub_key] = []

                        # Lisää vain jos ei ole jo listassa (vältetään duplikaatit)
                        entry = {
                            'content': value,
                            'source': doc_name,
                            'confidence': confidence
                        }

                        # Tarkista duplikaatit
                        is_duplicate = False
                        if isinstance(section_content[sub_key], list):
                            for existing in section_content[sub_key]:
                                if isinstance(existing, dict) and existing.get('content') == value:
                                    is_duplicate = True
                                    break

                        if not is_duplicate:
                            section_content[sub_key].append(entry)

            # Lisää osio vain jos siinä on sisältöä
            if section_content:
                result[section_key] = section_content

        return result

    def add_metadata(self, result, all_sections_data):
        """Lisää metadatan dokumenttiin"""
        metadata = OrderedDict()

        # Kerää company_name, jurisdiction, jne.
        for doc_name, sections in all_sections_data.items():
            if 'document_metadata' in sections and sections['document_metadata']:
                meta = sections['document_metadata']
                for key, value_dict in meta.items():
                    if isinstance(value_dict, dict):
                        value = value_dict.get('value', '')
                        if value and key not in metadata:
                            metadata[key] = value

        # Lisää sources_reviewed
        metadata['sources_reviewed'] = list(all_sections_data.keys())
        metadata['last_updated'] = datetime.now().isoformat()

        # Lisää metadata ensimmäiseksi
        result_with_metadata = OrderedDict()
        result_with_metadata['document_metadata'] = metadata
        result_with_metadata.update(result)

        return result_with_metadata

    def create_docx(self, json_data):
        """Luo DOCX-dokumentin JSON-datasta"""
        doc = Document()

        # Otsikko
        title = doc.add_heading('Due Diligence Report - Consolidated', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadata
        if 'document_metadata' in json_data:
            doc.add_heading('Document Metadata', 1)
            metadata = json_data['document_metadata']
            for key, value in metadata.items():
                p = doc.add_paragraph()
                p.add_run(f"{key}: ").bold = True
                p.add_run(str(value))

        # Käy läpi kaikki sektiot
        section_titles = {
            '1_executive_summary': '1. Executive Summary',
            '2_company_overview': '2. Company Overview',
            '3_market_and_go_to_market': '3. Market and Go-to-Market',
            '4_product_and_technology': '4. Product and Technology',
            '5_intellectual_property': '5. Intellectual Property',
            '6_team_and_organization': '6. Team and Organization',
            '7_customers_and_traction': '7. Customers and Traction',
            '8_operations_and_compliance': '8. Operations and Compliance',
            '9_financials': '9. Financials',
            '10_funding_and_capital_structure': '10. Funding and Capital Structure',
            '11_risks_and_dependencies': '11. Risks and Dependencies',
            '12_milestones_and_value_creation': '12. Milestones and Value Creation',
            '13_exit_considerations': '13. Exit Considerations',
            '14_open_questions_and_gaps': '14. Open Questions and Gaps',
            '15_appendices': '15. Appendices'
        }

        def process_section(data, parent_title='', level=1):
            """Rekursiivinen funktio sektioiden käsittelyyn"""
            if isinstance(data, dict):
                for key, value in data.items():
                    if key == 'document_metadata':
                        continue

                    # Tarkista onko tämä päätason sektio
                    section_title = section_titles.get(key, key.replace('_', ' ').title())

                    if isinstance(value, dict) and not all(k in ['content', 'source', 'confidence'] for k in value.keys()):
                        doc.add_heading(section_title, level)
                        process_section(value, section_title, level + 1)
                    elif isinstance(value, list):
                        doc.add_heading(section_title, level)
                        for item in value:
                            if isinstance(item, dict) and 'content' in item:
                                content = item.get('content', '')
                                if isinstance(content, str):
                                    p = doc.add_paragraph(content, style='List Bullet')
                                    if 'source' in item:
                                        source_run = p.add_run(f" [Source: {item['source']}]")
                                        source_run.font.size = Pt(8)
                                        source_run.font.color.rgb = RGBColor(128, 128, 128)
                            elif isinstance(item, str):
                                doc.add_paragraph(item, style='List Bullet')
                    elif isinstance(value, str):
                        doc.add_heading(section_title, level)
                        doc.add_paragraph(value)

        process_section(json_data)

        # Tallenna dokumentti
        doc.save(self.output_docx)
        print(f"\nDOCX-dokumentti tallennettu: {self.output_docx}")

    def process_all_documents(self):
        """Pääfunktio, joka käsittelee kaikki dokumentit"""
        print("Aloitetaan dokumenttien kasittely...")
        print("HUOM: Tama versio tekee useita API-kutsuja per dokumentti tarkkuuden vuoksi.\n")

        # Lataa template
        template = self.load_template()
        print(f"Template ladattu: {self.template_path}")

        # Etsi kaikki dokumentit
        document_files = list(self.docs_dir.glob('**/*'))
        document_files = [f for f in document_files if f.is_file() and f.suffix.lower() in ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt']]

        print(f"\nLoydettiin {len(document_files)} dokumenttia kasiteltavaksi")

        # Käsittele jokainen dokumentti
        all_sections_data = OrderedDict()

        for doc_file in document_files:
            print(f"\n{'='*60}")
            print(f"Kasitellaan: {doc_file.name}")
            print('='*60)

            # Lue dokumentti
            doc_text = self.read_document(doc_file)

            if not doc_text.strip():
                print(f"Dokumentti on tyhja tai lukeminen epaonnistui: {doc_file.name}")
                continue

            print(f"Luettu {len(doc_text)} merkkia")

            # Käsittele jokainen osio erikseen
            doc_sections = OrderedDict()

            for section_key, section_template in template.items():
                print(f"  Analysoidaan osio: {section_key}...", end=' ')

                extracted = self.extract_section_from_document(
                    doc_text,
                    doc_file.name,
                    section_key,
                    section_template
                )

                if extracted:
                    doc_sections[section_key] = extracted
                    print(f"[OK - {len(extracted)} kohtaa]")
                else:
                    doc_sections[section_key] = None
                    print("[Ei tietoa]")

            all_sections_data[doc_file.name] = doc_sections

        # Rakenna konsolidoitu dokumentti
        print("\n" + "="*60)
        print("Rakennetaan konsolidoitua dokumenttia...")
        print("="*60)

        consolidated = self.build_consolidated_document(all_sections_data, template)
        consolidated = self.add_metadata(consolidated, all_sections_data)

        # Tallenna JSON
        with open(self.output_json, 'w', encoding='utf-8') as f:
            json.dump(consolidated, f, indent=2, ensure_ascii=False)
        print(f"JSON tallennettu: {self.output_json}")

        # Tulosta tilastot
        print(f"\nLoydetyt osiot ({len([k for k in consolidated.keys() if k != 'document_metadata'])}):")
        for key in consolidated.keys():
            if key != 'document_metadata':
                print(f"  - {key}")

        # Luo DOCX
        print("\nLuodaan DOCX-dokumenttia...")
        self.create_docx(consolidated)

        print("\nValmis! Dokumentit kasitelty onnistuneesti.")
        print(f"  - JSON: {self.output_json}")
        print(f"  - DOCX: {self.output_docx}")


if __name__ == "__main__":
    processor = DocumentProcessor()
    processor.process_all_documents()
