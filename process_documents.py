import os
import json
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
import PyPDF2
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from collections import OrderedDict

# Lataa ympäristömuuttujat
load_dotenv()

class DocumentProcessor:
    def __init__(self):
        self.client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
        self.docs_dir = Path('docs')
        self.template_path = Path('master-document-template.json')
        self.output_json = Path('consolidated_due_diligence.json')
        self.output_docx = Path('consolidated_due_diligence.docx')

    def read_pdf(self, file_path):
        """Lukee PDF-tiedoston ja palauttaa tekstin"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"Virhe PDF:n lukemisessa {file_path}: {e}")
            return ""

    def read_excel(self, file_path):
        """Lukee Excel-tiedoston ja palauttaa tekstimuotoisen esityksen"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n=== Taulukko: {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            return text
        except Exception as e:
            print(f"Virhe Excel:in lukemisessa {file_path}: {e}")
            return ""

    def read_pptx(self, file_path):
        """Lukee PowerPoint-tiedoston ja palauttaa tekstin"""
        try:
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== Dia {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Virhe PowerPoint:in lukemisessa {file_path}: {e}")
            return ""

    def read_document(self, file_path):
        """Lukee dokumentin sen tyypin mukaan"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension == '.pdf':
            return self.read_pdf(file_path)
        elif extension in ['.xlsx', '.xls']:
            return self.read_excel(file_path)
        elif extension in ['.pptx', '.ppt']:
            return self.read_pptx(file_path)
        else:
            print(f"Tuntematon tiedostotyyppi: {extension}")
            return ""

    def load_template(self):
        """Lataa master template"""
        with open(self.template_path, 'r', encoding='utf-8') as f:
            return json.load(f, object_pairs_hook=OrderedDict)

    def create_template_summary(self, template):
        """Luo tiiviin yhteenvedon templatesta AI:lle"""
        summary = ""
        for section_key, section_data in template.items():
            if not isinstance(section_data, dict):
                continue

            summary += f"\n{section_key}:\n"

            for sub_key, sub_data in section_data.items():
                if isinstance(sub_data, dict):
                    instruction = sub_data.get('instruction', '')
                    update_rule = sub_data.get('update_rule', '')

                    if update_rule != 'locked' and instruction:
                        summary += f"  {sub_key}: {instruction}\n"

        return summary

    def extract_all_from_document(self, document_text, document_name, template):
        """Poimii KAIKKI tiedot dokumentista yhdellä API-kutsulla"""
        print(f"\nAnalysoidaan: {document_name}")

        template_summary = self.create_template_summary(template)

        prompt = f"""Olet tarkka Due Diligence -dokumenttianalyytikko. Analysoi tämä dokumentti ja poimii KAIKKI relevantti tieto rakenteelliseen muotoon.

KRIITTISET SÄÄNNÖT:
1. ÄLÄ KEKSI mitään - palauta VAIN eksplisiittisesti dokumentissa mainittu tieto
2. ÄLÄ HALLUSINOI - jos et ole varma, älä sisällytä
3. Poimii KAIKKI relevantti tieto - älä jätä mitään pois
4. Palauta tieto hierarkisessa JSON-muodossa joka vastaa template-rakennetta

Dokumentin nimi: {document_name}

TEMPLATE-RAKENNE JA OHJEET:
{template_summary}

DOKUMENTIN SISÄLTÖ:
{document_text}

Palauta JSON-objekti joka noudattaa template-rakennetta. Sisällytä VAIN ne osiot ja alakohdat joihin dokumentista löytyy tietoa.

Muoto:
{{
  "section_key": {{
    "subsection_key": {{
      "value": "Poimittu tieto dokumentista",
      "confidence": "high/medium/low"
    }}
  }}
}}

Jos osiossa on useita tietoja, palauta lista:
{{
  "section_key": {{
    "subsection_key": [
      {{"value": "Tieto 1", "confidence": "high"}},
      {{"value": "Tieto 2", "confidence": "medium"}}
    ]
  }}
}}

Palauta VAIN JSON. Älä lisää mitään selitystä."""

        try:
            print(f"  Lahetetaan API-kutsu... ({len(document_text)} merkkia)")

            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "Olet tarkka dokumenttianalyytikko. Palauta VAIN validi JSON, ei muuta tekstiä. Poimii KAIKKI relevantti tieto dokumentista."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.05,
                max_tokens=8192
            )

            result = response.choices[0].message.content.strip()

            # Etsi JSON-lohko
            if "```json" in result:
                result = result.split("```json")[1].split("```")[0].strip()
            elif "```" in result:
                result = result.split("```")[1].split("```")[0].strip()

            try:
                extracted = json.loads(result, object_pairs_hook=OrderedDict)

                # Laske tietopisteet
                count = 0
                for section_data in extracted.values():
                    if isinstance(section_data, dict):
                        for item in section_data.values():
                            if isinstance(item, list):
                                count += len(item)
                            else:
                                count += 1

                print(f"  Poimittu {count} tietopistetta")
                return extracted

            except json.JSONDecodeError as e:
                print(f"  JSON-parsinta epaonnistui: {e}")
                print(f"  Vastaus: {result[:500]}")
                return {}

        except Exception as e:
            print(f"  Virhe API-kutsussa: {e}")
            return {}

    def merge_document_data(self, all_documents_data, template):
        """Yhdistää kaikkien dokumenttien tiedot säilyttäen järjestyksen"""
        result = OrderedDict()

        # Käy läpi template järjestyksessä
        for section_key in template.keys():
            if section_key == 'document_metadata':
                continue

            if not isinstance(template[section_key], dict):
                continue

            section_content = OrderedDict()

            # Kerää tiedot kaikista dokumenteista
            for doc_name, doc_data in all_documents_data.items():
                if section_key not in doc_data:
                    continue

                section_data = doc_data[section_key]

                if not isinstance(section_data, dict):
                    continue

                # Käsittele alakohdat
                for sub_key, sub_value in section_data.items():
                    # Hae update_rule templatesta
                    update_rule = 'append'
                    if sub_key in template[section_key] and isinstance(template[section_key][sub_key], dict):
                        update_rule = template[section_key][sub_key].get('update_rule', 'append')

                    if update_rule == 'locked':
                        continue

                    # Käsittele arvot
                    if isinstance(sub_value, list):
                        # Lista tietoja
                        if update_rule == 'append':
                            if sub_key not in section_content:
                                section_content[sub_key] = []
                            elif not isinstance(section_content[sub_key], list):
                                # Muunna str -> list
                                section_content[sub_key] = [section_content[sub_key]]

                            for item in sub_value:
                                if isinstance(item, dict):
                                    value = item.get('value', '')
                                    confidence = item.get('confidence', 'medium')

                                    if value:
                                        entry = {
                                            'content': value,
                                            'source': doc_name,
                                            'confidence': confidence
                                        }

                                        # Vältetään duplikaatit
                                        is_dup = False
                                        for existing in section_content[sub_key]:
                                            if isinstance(existing, dict) and existing.get('content') == value:
                                                is_dup = True
                                                break

                                        if not is_dup:
                                            section_content[sub_key].append(entry)
                        else:
                            # overwrite
                            section_content[sub_key] = sub_value

                    elif isinstance(sub_value, dict):
                        # Yksittäinen tieto
                        value = sub_value.get('value', '')

                        if value:
                            if update_rule == 'overwrite':
                                section_content[sub_key] = value
                            elif update_rule == 'append':
                                if sub_key not in section_content:
                                    section_content[sub_key] = []
                                elif not isinstance(section_content[sub_key], list):
                                    # Muunna str -> list
                                    old_val = section_content[sub_key]
                                    section_content[sub_key] = [old_val]

                                entry = {
                                    'content': value,
                                    'source': doc_name,
                                    'confidence': sub_value.get('confidence', 'medium')
                                }

                                is_dup = False
                                for existing in section_content[sub_key]:
                                    if isinstance(existing, dict) and existing.get('content') == value:
                                        is_dup = True
                                        break

                                if not is_dup:
                                    section_content[sub_key].append(entry)

            # Lisää osio jos siinä on sisältöä
            if section_content:
                result[section_key] = section_content

        return result

    def add_metadata(self, result, all_documents_data):
        """Lisää metadatan"""
        metadata = OrderedDict()

        # Kerää metadata-kentät
        for doc_name, doc_data in all_documents_data.items():
            if 'document_metadata' in doc_data:
                meta = doc_data['document_metadata']
                if isinstance(meta, dict):
                    for key, value_data in meta.items():
                        if isinstance(value_data, dict):
                            value = value_data.get('value', '')
                            if value and key not in metadata:
                                metadata[key] = value

        metadata['sources_reviewed'] = list(all_documents_data.keys())
        metadata['last_updated'] = datetime.now().isoformat()

        result_with_metadata = OrderedDict()
        result_with_metadata['document_metadata'] = metadata
        result_with_metadata.update(result)

        return result_with_metadata

    def create_docx(self, json_data):
        """Luo DOCX-dokumentin"""
        doc = Document()

        title = doc.add_heading('Due Diligence Report - Consolidated', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        if 'document_metadata' in json_data:
            doc.add_heading('Document Metadata', 1)
            metadata = json_data['document_metadata']
            for key, value in metadata.items():
                p = doc.add_paragraph()
                p.add_run(f"{key}: ").bold = True
                p.add_run(str(value))

        section_titles = {
            '1_executive_summary': '1. Executive Summary',
            '2_company_overview': '2. Company Overview',
            '3_market_and_go_to_market': '3. Market and Go-to-Market',
            '4_product_and_technology': '4. Product and Technology',
            '5_intellectual_property': '5. Intellectual Property',
            '6_team_and_organization': '6. Team and Organization',
            '7_customers_and_traction': '7. Customers and Traction',
            '8_operations_and_compliance': '8. Operations and Compliance',
            '9_financials': '9. Financials',
            '10_funding_and_capital_structure': '10. Funding and Capital Structure',
            '11_risks_and_dependencies': '11. Risks and Dependencies',
            '12_milestones_and_value_creation': '12. Milestones and Value Creation',
            '13_exit_considerations': '13. Exit Considerations',
            '14_open_questions_and_gaps': '14. Open Questions and Gaps',
            '15_appendices': '15. Appendices'
        }

        def process_section(data, level=1):
            if isinstance(data, dict):
                for key, value in data.items():
                    if key == 'document_metadata':
                        continue

                    section_title = section_titles.get(key, key.replace('_', ' ').title())

                    if isinstance(value, dict) and not all(k in ['content', 'source', 'confidence'] for k in value.keys()):
                        doc.add_heading(section_title, level)
                        process_section(value, level + 1)
                    elif isinstance(value, list):
                        doc.add_heading(section_title, level)
                        for item in value:
                            if isinstance(item, dict) and 'content' in item:
                                content = item.get('content', '')
                                if isinstance(content, str):
                                    p = doc.add_paragraph(content, style='List Bullet')
                                    if 'source' in item:
                                        source_run = p.add_run(f" [Source: {item['source']}]")
                                        source_run.font.size = Pt(8)
                                        source_run.font.color.rgb = RGBColor(128, 128, 128)
                            elif isinstance(item, str):
                                doc.add_paragraph(item, style='List Bullet')
                    elif isinstance(value, str):
                        doc.add_heading(section_title, level)
                        doc.add_paragraph(value)

        process_section(json_data)

        doc.save(self.output_docx)
        print(f"\nDOCX-dokumentti tallennettu: {self.output_docx}")

    def process_all_documents(self):
        """Pääfunktio"""
        print("Aloitetaan dokumenttien kasittely...")
        print("Versio 3: Yksi API-kutsu per dokumentti, kattava poiminta\n")

        template = self.load_template()
        print(f"Template ladattu: {self.template_path}")

        document_files = list(self.docs_dir.glob('**/*'))
        document_files = [f for f in document_files if f.is_file() and f.suffix.lower() in ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt']]

        print(f"Loydettiin {len(document_files)} dokumenttia\n")

        all_documents_data = OrderedDict()

        for doc_file in document_files:
            print(f"{'='*60}")

            doc_text = self.read_document(doc_file)

            if not doc_text.strip():
                print(f"Dokumentti tyhja: {doc_file.name}")
                continue

            print(f"Luettu {len(doc_text)} merkkia")

            extracted = self.extract_all_from_document(doc_text, doc_file.name, template)

            if extracted:
                all_documents_data[doc_file.name] = extracted

        print(f"\n{'='*60}")
        print("Yhdistetaan kaikkien dokumenttien tiedot...")
        print('='*60)

        consolidated = self.merge_document_data(all_documents_data, template)
        consolidated = self.add_metadata(consolidated, all_documents_data)

        with open(self.output_json, 'w', encoding='utf-8') as f:
            json.dump(consolidated, f, indent=2, ensure_ascii=False)
        print(f"JSON tallennettu: {self.output_json}")

        print(f"\nLoydetyt osiot ({len([k for k in consolidated.keys() if k != 'document_metadata'])}):")
        for key in consolidated.keys():
            if key != 'document_metadata':
                print(f"  - {key}")

        print("\nLuodaan DOCX...")
        self.create_docx(consolidated)

        print("\nValmis!")
        print(f"  - JSON: {self.output_json}")
        print(f"  - DOCX: {self.output_docx}")


if __name__ == "__main__":
    processor = DocumentProcessor()
    processor.process_all_documents()
