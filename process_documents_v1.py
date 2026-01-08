import os
import json
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
import PyPDF2
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import io

# Lataa ympäristömuuttujat
load_dotenv()

class DocumentProcessor:
    def __init__(self):
        self.client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
        self.docs_dir = Path('docs')
        self.template_path = Path('master-document-template.json')
        self.output_json = Path('consolidated_due_diligence.json')
        self.output_docx = Path('consolidated_due_diligence.docx')

    def read_pdf(self, file_path):
        """Lukee PDF-tiedoston ja palauttaa tekstin"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"Virhe PDF:n lukemisessa {file_path}: {e}")
            return ""

    def read_excel(self, file_path):
        """Lukee Excel-tiedoston ja palauttaa tekstimuotoisen esityksen"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n=== Taulukko: {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            return text
        except Exception as e:
            print(f"Virhe Excel:in lukemisessa {file_path}: {e}")
            return ""

    def read_pptx(self, file_path):
        """Lukee PowerPoint-tiedoston ja palauttaa tekstin"""
        try:
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== Dia {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Virhe PowerPoint:in lukemisessa {file_path}: {e}")
            return ""

    def read_document(self, file_path):
        """Lukee dokumentin sen tyypin mukaan"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension == '.pdf':
            return self.read_pdf(file_path)
        elif extension in ['.xlsx', '.xls']:
            return self.read_excel(file_path)
        elif extension in ['.pptx', '.ppt']:
            return self.read_pptx(file_path)
        else:
            print(f"Tuntematon tiedostotyyppi: {extension}")
            return ""

    def load_template(self):
        """Lataa master template"""
        with open(self.template_path, 'r', encoding='utf-8') as f:
            return json.load(f)

    def extract_data_from_document(self, document_text, document_name, template):
        """Käyttää OpenAI API:a tiedon poimintaan dokumentista"""
        print(f"\nAnalysoidaan dokumentti: {document_name}")

        # Rakenna prompt
        prompt = f"""Olet tarkka ja rehellinen dokumenttianalyytikko. Sinun tehtäväsi on analysoida Due Diligence -dokumentti ja poimia siitä vain faktat, jotka on eksplisiittisesti mainittu dokumentissa.

TÄRKEÄÄ:
- ÄLÄ KEKSI MITÄÄN
- ÄLÄ HALLUSINOI
- Palauta vain ne kohdat, joihin löytyy selkeää, todennettavaa tietoa dokumentista
- Jos tieto ei ole dokumentissa, älä sisällytä kyseistä kenttää vastaukseen
- Noudata tarkkaan kunkin kentän "instruction"-ohjeita

Dokumentin nimi: {document_name}

Template-rakenne ja ohjeet:
{json.dumps(template, indent=2, ensure_ascii=False)}

Analysoitava dokumentti:
{document_text[:30000]}

Palauta JSON-objekti, joka sisältää VAIN ne kentät ja osa-alueet, joihin dokumentista löytyy relevanttia tietoa. Älä sisällytä tyhjiä kenttiä tai kenttiä, joihin ei ole tietoa.

Palauta vastaus muodossa:
{{
  "section_path": "esim. 2_company_overview.2_1_company_description",
  "value": "poimittu tieto dokumentista",
  "confidence": "high/medium/low",
  "source_quote": "lyhyt lainaus dokumentista"
}}

Palauta lista kaikista löydetyistä tiedoista JSON-array muodossa."""

        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",  # Paras saatavilla oleva malli tarkkaan analyysiin
                messages=[
                    {"role": "system", "content": "Olet tarkka dokumenttianalyytikko, joka poimii vain faktoja dokumenteista ilman hallusinaatioita."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,  # Matala lämpötila vähentää hallusinaatioita
                max_tokens=4096
            )

            result = response.choices[0].message.content

            # Yritä parsia JSON-vastaus
            try:
                # Etsi JSON-lohko vastauksesta
                if "```json" in result:
                    result = result.split("```json")[1].split("```")[0].strip()
                elif "```" in result:
                    result = result.split("```")[1].split("```")[0].strip()

                extracted_data = json.loads(result)
                return extracted_data
            except json.JSONDecodeError as e:
                print(f"JSON-parsinta epäonnistui dokumentille {document_name}: {e}")
                print(f"Vastaus: {result[:500]}")
                return []

        except Exception as e:
            print(f"Virhe OpenAI API-kutsussa: {e}")
            return []

    def build_consolidated_document(self, all_extracted_data, template):
        """Rakentaa konsolidoidun dokumentin kaikesta poimitusta datasta"""
        result = {}
        sources = set()

        # Käy läpi jokainen dokumentti ja sen poiminta
        for doc_name, extracted_items in all_extracted_data.items():
            sources.add(doc_name)

            if not isinstance(extracted_items, list):
                continue

            for item in extracted_items:
                if not isinstance(item, dict):
                    continue

                section_path = item.get('section_path', '')
                value = item.get('value', '')

                if not section_path or not value:
                    continue

                # Rakenna polku hierarkiaan
                path_parts = section_path.split('.')
                current = result

                # Luo polku hierarkiaan
                for i, part in enumerate(path_parts[:-1]):
                    if part not in current:
                        current[part] = {}
                    current = current[part]

                # Lisää arvo
                last_key = path_parts[-1]

                # Hae update_rule templatesta
                template_current = template
                for part in path_parts:
                    if isinstance(template_current, dict) and part in template_current:
                        template_current = template_current[part]

                update_rule = template_current.get('update_rule', 'append') if isinstance(template_current, dict) else 'append'

                # Älä lisää locked-kenttiä
                if update_rule == 'locked':
                    continue

                # Käsittele update_rule
                if update_rule == 'overwrite':
                    current[last_key] = value
                elif update_rule == 'append':
                    if last_key not in current:
                        current[last_key] = []
                    if isinstance(current[last_key], list):
                        current[last_key].append({
                            'content': value,
                            'source': doc_name,
                            'confidence': item.get('confidence', 'medium')
                        })
                    else:
                        # Jos ei ole lista, tee siitä lista
                        old_value = current[last_key]
                        current[last_key] = [old_value, {
                            'content': value,
                            'source': doc_name,
                            'confidence': item.get('confidence', 'medium')
                        }]

        # Lisää metadata
        if 'document_metadata' not in result:
            result['document_metadata'] = {}

        result['document_metadata']['sources_reviewed'] = list(sources)
        result['document_metadata']['last_updated'] = datetime.now().isoformat()

        return result

    def create_docx(self, json_data):
        """Luo DOCX-dokumentin JSON-datasta"""
        doc = Document()

        # Otsikko
        title = doc.add_heading('Due Diligence Report - Consolidated', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadata
        if 'document_metadata' in json_data:
            doc.add_heading('Document Metadata', 1)
            metadata = json_data['document_metadata']
            for key, value in metadata.items():
                p = doc.add_paragraph()
                p.add_run(f"{key}: ").bold = True
                p.add_run(str(value))

        # Käy läpi kaikki sektiot
        section_titles = {
            '1_executive_summary': '1. Executive Summary',
            '2_company_overview': '2. Company Overview',
            '3_market_and_go_to_market': '3. Market and Go-to-Market',
            '4_product_and_technology': '4. Product and Technology',
            '5_intellectual_property': '5. Intellectual Property',
            '6_team_and_organization': '6. Team and Organization',
            '7_customers_and_traction': '7. Customers and Traction',
            '8_operations_and_compliance': '8. Operations and Compliance',
            '9_financials': '9. Financials',
            '10_funding_and_capital_structure': '10. Funding and Capital Structure',
            '11_risks_and_dependencies': '11. Risks and Dependencies',
            '12_milestones_and_value_creation': '12. Milestones and Value Creation',
            '13_exit_considerations': '13. Exit Considerations',
            '14_open_questions_and_gaps': '14. Open Questions and Gaps',
            '15_appendices': '15. Appendices'
        }

        def process_section(data, parent_title='', level=1):
            """Rekursiivinen funktio sektioiden käsittelyyn"""
            if isinstance(data, dict):
                for key, value in data.items():
                    if key == 'document_metadata':
                        continue

                    # Tarkista onko tämä päätason sektio
                    section_title = section_titles.get(key, key.replace('_', ' ').title())

                    if isinstance(value, dict):
                        doc.add_heading(section_title, level)
                        process_section(value, section_title, level + 1)
                    elif isinstance(value, list):
                        doc.add_heading(section_title, level)
                        for item in value:
                            if isinstance(item, dict):
                                p = doc.add_paragraph(item.get('content', str(item)), style='List Bullet')
                                if 'source' in item:
                                    source_run = p.add_run(f" [Source: {item['source']}]")
                                    source_run.font.size = Pt(8)
                                    source_run.font.color.rgb = RGBColor(128, 128, 128)
                            else:
                                doc.add_paragraph(str(item), style='List Bullet')
                    else:
                        doc.add_heading(section_title, level)
                        doc.add_paragraph(str(value))
            elif isinstance(data, list):
                for item in data:
                    if isinstance(item, dict):
                        p = doc.add_paragraph(item.get('content', str(item)), style='List Bullet')
                        if 'source' in item:
                            source_run = p.add_run(f" [Source: {item['source']}]")
                            source_run.font.size = Pt(8)
                            source_run.font.color.rgb = RGBColor(128, 128, 128)
                    else:
                        doc.add_paragraph(str(item), style='List Bullet')

        process_section(json_data)

        # Tallenna dokumentti
        doc.save(self.output_docx)
        print(f"\nDOCX-dokumentti tallennettu: {self.output_docx}")

    def process_all_documents(self):
        """Pääfunktio, joka käsittelee kaikki dokumentit"""
        print("Aloitetaan dokumenttien käsittely...")

        # Lataa template
        template = self.load_template()
        print(f"Template ladattu: {self.template_path}")

        # Etsi kaikki dokumentit
        document_files = list(self.docs_dir.glob('**/*'))
        document_files = [f for f in document_files if f.is_file() and f.suffix.lower() in ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt']]

        print(f"\nLöydettiin {len(document_files)} dokumenttia käsiteltäväksi")

        # Käsittele jokainen dokumentti
        all_extracted_data = {}

        for doc_file in document_files:
            print(f"\nKäsitellään: {doc_file.name}")

            # Lue dokumentti
            doc_text = self.read_document(doc_file)

            if not doc_text.strip():
                print(f"Dokumentti on tyhjä tai lukeminen epäonnistui: {doc_file.name}")
                continue

            print(f"Luettu {len(doc_text)} merkkiä dokumentista {doc_file.name}")

            # Poimii tiedot OpenAI:n avulla
            extracted = self.extract_data_from_document(doc_text, doc_file.name, template)
            all_extracted_data[doc_file.name] = extracted

            print(f"Poimittu {len(extracted) if isinstance(extracted, list) else 0} datapistettä")

        # Rakenna konsolidoitu dokumentti
        print("\nRakennetaan konsolidoitua dokumenttia...")
        consolidated = self.build_consolidated_document(all_extracted_data, template)

        # Tallenna JSON
        with open(self.output_json, 'w', encoding='utf-8') as f:
            json.dump(consolidated, f, indent=2, ensure_ascii=False)
        print(f"JSON tallennettu: {self.output_json}")

        # Luo DOCX
        print("\nLuodaan DOCX-dokumenttia...")
        self.create_docx(consolidated)

        print("\nValmis! Dokumentit kasitelty onnistuneesti.")
        print(f"  - JSON: {self.output_json}")
        print(f"  - DOCX: {self.output_docx}")


if __name__ == "__main__":
    processor = DocumentProcessor()
    processor.process_all_documents()
