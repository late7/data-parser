from __future__ import annotations

import argparse
import copy
import json
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from dotenv import load_dotenv


@dataclass(frozen=True)
class TextChunk:
    source_file: str
    location: str
    text: str


def _resolve_model_name(requested: str) -> str:
    """Resolve user-facing model alias to an API model id.

    Some environments refer to models as "chatgpt-*" while the OpenAI API expects
    "gpt-*" ids.
    """

    name = (requested or "").strip()
    if not name:
        return name
    if name.startswith("chatgpt-"):
        return "gpt-" + name[len("chatgpt-") :]
    return name


def _read_json(path: Path) -> Dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def _write_json(path: Path, data: Dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _is_leaf_section(node: Any) -> bool:
    return isinstance(node, dict) and "update_rule" in node and "instruction" in node


def _collect_leaf_pointers(template: Dict[str, Any]) -> List[Tuple[str, Dict[str, Any]]]:
    leaves: List[Tuple[str, Dict[str, Any]]] = []

    def walk(obj: Any, path_parts: List[str]) -> None:
        if _is_leaf_section(obj):
            pointer = "/" + "/".join(path_parts)
            leaves.append((pointer, obj))
            return
        if isinstance(obj, dict):
            for k, v in obj.items():
                walk(v, path_parts + [k])

    walk(template, [])
    return leaves


def _get_by_pointer(root: Dict[str, Any], pointer: str) -> Dict[str, Any]:
    parts = [p for p in pointer.split("/") if p]
    cur: Any = root
    for p in parts:
        if not isinstance(cur, dict) or p not in cur:
            raise KeyError(f"Pointer not found: {pointer}")
        cur = cur[p]
    if not isinstance(cur, dict):
        raise TypeError(f"Pointer does not resolve to dict: {pointer}")
    return cur


def _normalize_whitespace(s: str) -> str:
    s = s.replace("\u00a0", " ")
    s = re.sub(r"[\t\r]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    s = re.sub(r"[ ]{2,}", " ", s)
    return s.strip()


def _coerce_json_object(text: str) -> Dict[str, Any]:
    """Coerce model output into a JSON object.

    We ask the model for strict JSON, but this makes parsing resilient to
    accidental code fences or surrounding text.
    """

    s = (text or "").strip()
    if s.startswith("```"):
        # Remove common fenced blocks: ```json ... ```
        s = re.sub(r"^```[a-zA-Z0-9_-]*\n", "", s)
        s = re.sub(r"\n```$", "", s)
        s = s.strip()

    start = s.find("{")
    end = s.rfind("}")
    if start != -1 and end != -1 and end > start:
        s = s[start : end + 1]

    parsed = json.loads(s)
    if not isinstance(parsed, dict):
        raise ValueError("Model output JSON was not an object")
    return parsed


def _chunk_text(text: str, max_chars: int) -> List[str]:
    text = _normalize_whitespace(text)
    if len(text) <= max_chars:
        return [text]

    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        # Try to split on paragraph boundary.
        split = text.rfind("\n\n", start, end)
        if split == -1 or split <= start + max_chars // 2:
            split = end
        chunk = text[start:split].strip()
        if chunk:
            chunks.append(chunk)
        start = split
    return chunks


def extract_text_chunks_from_docs(docs_dir: Path, max_chunk_chars: int) -> List[TextChunk]:
    chunks: List[TextChunk] = []

    for path in sorted(docs_dir.iterdir()):
        if path.is_dir():
            continue

        suffix = path.suffix.lower()
        if suffix == ".pdf":
            chunks.extend(_extract_pdf(path, max_chunk_chars))
        elif suffix == ".pptx":
            chunks.extend(_extract_pptx(path, max_chunk_chars))
        elif suffix == ".xlsx":
            chunks.extend(_extract_xlsx(path, max_chunk_chars))
        elif suffix in {".txt", ".md", ".json"}:
            text = path.read_text(encoding="utf-8", errors="replace")
            for i, chunk in enumerate(_chunk_text(text, max_chunk_chars), start=1):
                chunks.append(TextChunk(path.name, f"chunk {i}", chunk))
        else:
            # Unsupported type; skip silently to avoid accidental hallucination.
            continue

    return chunks


def _extract_pdf(path: Path, max_chunk_chars: int) -> List[TextChunk]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    out: List[TextChunk] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = _normalize_whitespace(text)
        if not text:
            continue
        for i, chunk in enumerate(_chunk_text(text, max_chunk_chars), start=1):
            loc = f"page {page_index}" if i == 1 else f"page {page_index}, chunk {i}"
            out.append(TextChunk(path.name, loc, chunk))
    return out


def _extract_pptx(path: Path, max_chunk_chars: int) -> List[TextChunk]:
    from pptx import Presentation

    pres = Presentation(str(path))
    out: List[TextChunk] = []
    for slide_index, slide in enumerate(pres.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = getattr(shape, "text")
                if t:
                    parts.append(t)
        text = _normalize_whitespace("\n".join(parts))
        if not text:
            continue
        for i, chunk in enumerate(_chunk_text(text, max_chunk_chars), start=1):
            loc = f"slide {slide_index}" if i == 1 else f"slide {slide_index}, chunk {i}"
            out.append(TextChunk(path.name, loc, chunk))
    return out


def _extract_xlsx(path: Path, max_chunk_chars: int) -> List[TextChunk]:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), data_only=True)
    out: List[TextChunk] = []

    for sheet in wb.worksheets:
        rows: List[str] = []
        for row in sheet.iter_rows(values_only=True):
            # Keep rows that have any content
            if not row:
                continue
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                rows.append("\t".join(values))

        text = _normalize_whitespace("\n".join(rows))
        if not text:
            continue

        for i, chunk in enumerate(_chunk_text(text, max_chunk_chars), start=1):
            loc = f"sheet {sheet.title}" if i == 1 else f"sheet {sheet.title}, chunk {i}"
            out.append(TextChunk(path.name, loc, chunk))

    return out


def build_output_skeleton(template: Dict[str, Any]) -> Dict[str, Any]:
    out = copy.deepcopy(template)
    for pointer, leaf in _collect_leaf_pointers(out):
        update_rule = str(leaf.get("update_rule", "")).lower()
        if update_rule == "locked":
            continue
        # Add extraction slots. Keep instructions intact for traceability.
        leaf.setdefault("extracted", None)
        leaf.setdefault("evidence", [])
    return out


def _make_field_catalog(template: Dict[str, Any]) -> List[Dict[str, str]]:
    catalog: List[Dict[str, str]] = []
    for pointer, leaf in _collect_leaf_pointers(template):
        update_rule = str(leaf.get("update_rule", "")).lower()
        if update_rule == "locked":
            continue
        catalog.append(
            {
                "pointer": pointer,
                "update_rule": update_rule,
                "instruction": str(leaf.get("instruction", "")),
            }
        )
    return catalog


def call_openai_extract(
    *,
    client: Any,
    model: str,
    field_catalog: List[Dict[str, str]],
    chunk: TextChunk,
) -> List[Dict[str, Any]]:
    """Returns a list of extractions.

    Each extraction should contain:
      - pointer: JSON pointer into template
      - value: extracted factual content
      - evidence: list[{quote, location}]

    Validation against hallucination is done after the call.
    """

    system_prompt = (
        "You are an analyst extracting facts for a Due Diligence report. "
        "Only extract information that is explicitly present in the provided source text. "
        "Do not guess, infer, or use outside knowledge. "
        "Every extracted item MUST include one or more exact quotes copied verbatim from the source text. "
        "If you cannot support a claim with an exact quote, do not output it. "
        "Return strictly valid JSON only (no markdown)."
    )

    user_payload = {
        "task": "Extract only relevant info from the source for the specified master document fields.",
        "source": {
            "file": chunk.source_file,
            "location": chunk.location,
            "text": chunk.text,
        },
        "fields": field_catalog,
        "output_format": {
            "extractions": [
                {
                    "pointer": "/2_company_overview/2_1_company_description",
                    "value": "<factual text strictly supported by quotes>",
                    "evidence": [
                        {"quote": "<exact substring from source text>", "location": "<use given location>"}
                    ],
                }
            ]
        },
        "rules": [
            "Only include fields that have relevant info in the source text.",
            "Use the pointer values exactly as provided.",
            "Evidence quotes must be exact substrings from the source text.",
            "If nothing relevant exists, return {\"extractions\": []}.",
        ],
    }

    try:
        resp = client.responses.create(
            model=model,
            input=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": json.dumps(user_payload, ensure_ascii=False)},
            ],
        )
    except Exception as e:
        # Provide a clearer hint for the most common failure.
        msg = str(e)
        if "model_not_found" in msg or "does not exist" in msg:
            raise RuntimeError(
                f"OpenAI API model not found: '{model}'. "
                f"Try '--model gpt-4o-mini' (cheap) or set OPENAI_MODEL in .env, "
                f"or verify your account has access to the requested model."
            ) from e
        raise

    text = getattr(resp, "output_text", None)
    if not text:
        # Fallback: try to reconstruct from output content
        try:
            text = resp.output[0].content[0].text  # type: ignore[attr-defined]
        except Exception:
            raise RuntimeError("OpenAI response did not contain output_text")

    parsed = _coerce_json_object(text)
    extractions = parsed.get("extractions", [])
    if not isinstance(extractions, list):
        return []
    return extractions


def validate_and_merge(
    *,
    output_doc: Dict[str, Any],
    template: Dict[str, Any],
    chunk: TextChunk,
    extractions: List[Dict[str, Any]],
) -> int:
    """Validates extractions and merges into output_doc.

    Returns count of merged items.
    """

    merged = 0
    for item in extractions:
        if not isinstance(item, dict):
            continue
        pointer = item.get("pointer")
        value = item.get("value")
        evidence = item.get("evidence")

        if not isinstance(pointer, str) or not pointer.startswith("/"):
            continue

        # Ensure pointer exists in template and is a leaf
        try:
            template_leaf = _get_by_pointer(template, pointer)
            if not _is_leaf_section(template_leaf):
                continue
        except Exception:
            continue

        update_rule = str(template_leaf.get("update_rule", "")).lower()
        if update_rule == "locked":
            continue

        if not isinstance(evidence, list) or not evidence:
            continue

        # Validate evidence quotes are substrings
        valid_evidence: List[Dict[str, str]] = []
        for ev in evidence:
            if not isinstance(ev, dict):
                continue
            quote = ev.get("quote")
            if not isinstance(quote, str) or not quote.strip():
                continue
            if quote not in chunk.text:
                continue
            valid_evidence.append(
                {
                    "source_file": chunk.source_file,
                    "source_location": chunk.location,
                    "quote": quote.strip(),
                }
            )

        if not valid_evidence:
            continue

        out_leaf = _get_by_pointer(output_doc, pointer)
        out_leaf.setdefault("evidence", [])
        out_leaf["evidence"].extend(valid_evidence)

        # Merge extracted value according to update_rule
        if update_rule == "overwrite":
            if isinstance(value, str) and value.strip():
                out_leaf["extracted"] = value.strip()
                merged += 1
        elif update_rule == "append":
            if isinstance(value, str) and value.strip():
                if out_leaf.get("extracted") is None:
                    out_leaf["extracted"] = []
                if isinstance(out_leaf.get("extracted"), list):
                    out_leaf["extracted"].append(value.strip())
                    merged += 1
        else:
            # Default: append-like behavior
            if isinstance(value, str) and value.strip():
                if out_leaf.get("extracted") is None:
                    out_leaf["extracted"] = []
                if isinstance(out_leaf.get("extracted"), list):
                    out_leaf["extracted"].append(value.strip())
                    merged += 1

    return merged


def fill_sources_reviewed(output_doc: Dict[str, Any], source_files: List[str]) -> None:
    try:
        leaf = _get_by_pointer(output_doc, "/document_metadata/sources_reviewed")
    except Exception:
        return

    if str(leaf.get("update_rule", "")).lower() == "locked":
        return

    leaf.setdefault("extracted", [])
    existing = leaf.get("extracted")
    existing_list: List[str] = existing if isinstance(existing, list) else []
    leaf["extracted"] = list(dict.fromkeys(existing_list + source_files))


def add_open_questions_for_missing(output_doc: Dict[str, Any], template: Dict[str, Any]) -> None:
    missing: List[str] = []
    for pointer, leaf in _collect_leaf_pointers(template):
        update_rule = str(leaf.get("update_rule", "")).lower()
        if update_rule == "locked":
            continue
        out_leaf = _get_by_pointer(output_doc, pointer)
        extracted = out_leaf.get("extracted", None)
        if extracted is None or extracted == "" or extracted == []:
            missing.append(pointer)

    if not missing:
        return

    try:
        gaps_leaf = _get_by_pointer(output_doc, "/14_open_questions_and_gaps")
    except Exception:
        return

    if str(gaps_leaf.get("update_rule", "")).lower() == "locked":
        return

    gaps_leaf.setdefault("extracted", [])
    if gaps_leaf.get("extracted") is None:
        gaps_leaf["extracted"] = []

    if isinstance(gaps_leaf.get("extracted"), list):
        gaps_leaf["extracted"].append(
            "Missing or not evidenced in provided documents for fields: " + ", ".join(missing)
        )


def render_docx(output_doc: Dict[str, Any], docx_path: Path) -> None:
    from docx import Document

    def is_leaf(obj: Any) -> bool:
        return _is_leaf_section(obj)

    def walk(doc: Document, obj: Any, heading: str, level: int) -> None:
        if isinstance(obj, dict) and is_leaf(obj):
            extracted = obj.get("extracted")
            evidence = obj.get("evidence")
            if extracted is None or extracted == "" or extracted == []:
                return
            doc.add_heading(heading, level=level)
            if isinstance(extracted, list):
                for item in extracted:
                    doc.add_paragraph(str(item), style="List Bullet")
            else:
                doc.add_paragraph(str(extracted))

            if isinstance(evidence, list) and evidence:
                doc.add_paragraph("Evidence:")
                for ev in evidence:
                    try:
                        src = f"{ev.get('source_file')} ({ev.get('source_location')})"
                        quote = ev.get("quote")
                        doc.add_paragraph(f"{src}: {quote}", style="List Bullet")
                    except Exception:
                        continue
            return

        if isinstance(obj, dict):
            for k, v in obj.items():
                next_heading = f"{heading} / {k}" if heading else str(k)
                next_level = min(9, level + 1) if heading else 1
                walk(doc, v, next_heading, next_level)

    doc = Document()
    doc.add_heading("Due Diligence Master Document (Draft)", level=0)
    doc.add_paragraph(
        "This draft is automatically compiled from provided source documents. "
        "It includes only items supported by direct quotes (see evidence bullets)."
    )

    walk(doc, output_doc, "", 1)

    docx_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(docx_path))


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Compile a Due Diligence master JSON from docs/ using OpenAI Responses API, and render DOCX for review."
    )
    parser.add_argument("--docs", default="docs", help="Docs directory")
    parser.add_argument(
        "--template",
        default="master-document-template.json",
        help="Template JSON path",
    )
    parser.add_argument("--out-json", default="output/master_document.json")
    parser.add_argument("--out-docx", default="output/master_document.docx")
    parser.add_argument("--model", default=os.getenv("OPENAI_MODEL", "chatgpt-5.2"))
    parser.add_argument("--max-chunk-chars", type=int, default=12000)
    args = parser.parse_args()
    model = _resolve_model_name(args.model)

    load_dotenv()
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY missing (expected in .env)")

    from openai import OpenAI

    client = OpenAI(api_key=api_key)

    docs_dir = Path(args.docs)
    template_path = Path(args.template)
    out_json = Path(args.out_json)
    out_docx = Path(args.out_docx)

    template = _read_json(template_path)
    output_doc = build_output_skeleton(template)

    field_catalog = _make_field_catalog(template)

    text_chunks = extract_text_chunks_from_docs(docs_dir, args.max_chunk_chars)
    source_files = sorted({c.source_file for c in text_chunks})
    fill_sources_reviewed(output_doc, source_files)

    merged_total = 0
    for chunk in text_chunks:
        extractions = call_openai_extract(
            client=client,
            model=model,
            field_catalog=field_catalog,
            chunk=chunk,
        )
        merged_total += validate_and_merge(
            output_doc=output_doc,
            template=template,
            chunk=chunk,
            extractions=extractions,
        )

    add_open_questions_for_missing(output_doc, template)

    _write_json(out_json, output_doc)
    render_docx(output_doc, out_docx)

    print(f"Wrote JSON: {out_json}")
    print(f"Wrote DOCX: {out_docx}")
    print(f"Merged items: {merged_total}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
