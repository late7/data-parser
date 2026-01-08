import os
import json
import glob
from dataclasses import dataclass
from typing import Dict, Any, List
from dotenv import load_dotenv
from openai import OpenAI
import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import pypdf
import pptx
import openpyxl

# Load environment variables
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
MODEL_NAME = "chatgpt-5.2"  # As requested
DOCS_DIR = "docs"
TEMPLATE_FILE = "master-document-template.json"
OUTPUT_JSON = "master_document.json"
OUTPUT_DOCX = "Due_Diligence_Report.docx"

client = OpenAI(api_key=OPENAI_API_KEY)

def extract_text_from_pdf(filepath):
    text = ""
    try:
        reader = pypdf.PdfReader(filepath)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
    return text

def extract_text_from_pptx(filepath):
    text = ""
    try:
        prs = pptx.Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {filepath}: {e}")
    return text

def extract_text_from_xlsx(filepath):
    text = ""
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    text += " ".join(row_text) + "\n"
    except Exception as e:
        print(f"Error reading XLSX {filepath}: {e}")
    return text

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    elif ext in [".xlsx", ".xls"]:
         return extract_text_from_xlsx(filepath)
    else:
        # Assume text based
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file {filepath}: {e}")
            return ""

def process_file_with_ai(filename, content, template_structure):
    """
    Sends file content and template instructions to LLM to extract data.
    """
    # Create a simplified schema description for the prompt
    # We want the LLM to return a JSON that matches the template keys
    
    prompt = f"""
    You are an expert Due Diligence Analyst. 
    Analyze the following document: '{filename}'.
    
    Your task is to extract relevant information to populate a Due Diligence Report based on the provided JSON schema instructions.
    
    The Schema has specific sections with 'instruction' fields. 
    For each field in the schema, checks if the document contains information compliant with the 'instruction'.
    
    Return a JSON object that strictly follows the structure of the provided schema (keys matching).
    - If a field is found, provide the value.
    - If a field is NOT found or the document is not relevant to it, leave it null or empty string.
    - Do NOT invent information. Only extract what is explicitly stated.
    - Respect 'locked' fields: do not provide values for them (they are static).
    
    Template Schema (with instructions):
    {json.dumps(template_structure, indent=2)}
    
    Document Content:
    {content[:100000]} # Limit characters to avoid token limits if necessary, though newer models handle large context.
    """

    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        print(f"Error processing {filename} with AI: {e}")
        # Fallback to gpt-4o if chatgpt-5.2 doesn't exist
        if "model_not_found" in str(e) or "404" in str(e):
             print("Model not found, falling back to gpt-4o")
             try:
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[{"role": "user", "content": prompt}],
                    response_format={"type": "json_object"}
                )
                return json.loads(response.choices[0].message.content)
             except Exception as e2:
                 print(f"Fallback failed: {e2}")
                 return {}
        return {}

def merge_data(master_data, new_data, template):
    """
    Merges extracted data into master data based on rules in template.
    """
    for section_key, section_val in new_data.items():
        if section_key not in master_data:
            continue # Should be same structure
            
        if isinstance(section_val, dict):
            # Recurse for nested fields (subsection)
            if section_key not in master_data:
                 master_data[section_key] = {}
            
            for field_key, field_val in section_val.items():
                if not field_val:
                    continue # No data extracted
                
                # Get update rule from template
                # Template structure: Section -> SubSection -> {update_rule, instruction, value(initially empty)}
                # But logic: Master starts as finding keys. 
                # Let's look at template: "1_executive_summary": { "1_1_investment_snapshot": { "update_rule": ... } }
                
                rule = "overwrite" # Default
                if section_key in template and field_key in template[section_key]:
                     rule = template[section_key][field_key].get("update_rule", "overwrite")
                elif section_key == "document_metadata" and field_key in template["document_metadata"]:
                     rule = template[section_key][field_key].get("update_rule", "overwrite")

                if rule == "locked":
                    continue
                
                current_val = master_data[section_key].get(field_key, "")
                
                # We need to store the value. The master_data structure initially matches template.
                # But template has metadata dicts. We want master_data to hold the CONTENT.
                # So we should probably initialize master_data to hold content strings, or keep the dict structure?
                # The user wants "uuden .json documentin... samalla rakenteella".
                # "Lisää tähän... ne kohdat".
                # I will maintain the structure: Section -> Field -> "content": "..."
                # Or just Section -> Field: "content"
                
                # Let's clean master_data initialization first.
                pass

def initialize_master_data(template):
    data = {}
    for key, val in template.items():
        if isinstance(val, dict) and "update_rule" not in val:
             # It's a section
             data[key] = {}
             for subkey, subval in val.items():
                  if isinstance(subval, dict) and "instruction" in subval:
                      data[key][subkey] = "" # Initialize empty content
        elif isinstance(val, dict) and "instruction" in val:
             # Root level field (like 14_open_questions...) wait, checking template again
             # 14_... is "update_rule": "append". It is a field directly.
             data[key] = ""
    return data

def update_master_document(master_doc, extracted, template):
    """
    Updates the master document holding values.
    """
    # 1. Recursive update
    for key, val in extracted.items():
        if key not in master_doc: 
            continue
        
        if isinstance(val, dict):
             for subkey, subval in val.items():
                 if not subval: continue
                 if subkey not in master_doc[key]: continue
                 
                 # Ensure subval is string
                 if isinstance(subval, list):
                     subval = "; ".join([str(x) for x in subval])
                 elif not isinstance(subval, str):
                     subval = str(subval)
                 
                 # Determine rule
                 rule = "overwrite"
                 if key in template and subkey in template[key]:
                      rule = template[key][subkey].get("update_rule", "overwrite")
                 
                 existing = master_doc[key][subkey]
                 if isinstance(existing, list): # Should be string now, but safety check
                      existing = "; ".join([str(x) for x in existing])
                 
                 if rule == "overwrite":
                      master_doc[key][subkey] = subval
                 elif rule == "append":
                      if existing:
                           master_doc[key][subkey] = existing + "\n\n" + subval
                      else:
                           master_doc[key][subkey] = subval
        else:
             # Top level field (should be str)
             if not val: continue 
             
             if isinstance(val, list):
                 val = "; ".join([str(x) for x in val])
             elif not isinstance(val, str):
                 val = str(val)

             rule = template.get(key, {}).get("update_rule", "overwrite")
             existing = master_doc[key]
             
             if rule == "overwrite":
                  master_doc[key] = val
             elif rule == "append":
                   if existing:
                       master_doc[key] = existing + "\n\n" + val
                   else:
                       master_doc[key] = val
    return master_doc

def generate_docx(data, filename):
    doc = docx.Document()
    
    # Styles
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    title = doc.add_heading("Due Diligence Report", 0)
    title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    
    # Metadata
    if "document_metadata" in data:
        doc.add_heading("Metadata", level=1)
        for k, v in data["document_metadata"].items():
             p = doc.add_paragraph()
             p.add_run(f"{k.replace('_', ' ').title()}: ").bold = True
             p.add_run(str(v))
    
    # Sections
    sorted_keys = sorted(data.keys())
    for key in sorted_keys:
        if key == "document_metadata": continue
        
        # Format Heading
        heading_text = key.replace("_", " ").title()
        # Remove numbers for cleaner look if desired, or keep them. keeping them helps sorting.
        doc.add_heading(heading_text, level=1)
        
        val = data[key]
        if isinstance(val, dict):
             for subkey, subval in val.items():
                  subheading = subkey.replace("_", " ").title()
                  doc.add_heading(subheading, level=2)
                  if subval:
                      doc.add_paragraph(subval)
                  else:
                      doc.add_paragraph("[No information found]")
        elif isinstance(val, str):
             if val:
                  doc.add_paragraph(val)
             else:
                  doc.add_paragraph("[No information found]")
                  
    doc.save(filename)
    print(f"Saved DOCX to {filename}")

def main():
    if not os.path.exists(TEMPLATE_FILE):
        print("Template not found!")
        return

    with open(TEMPLATE_FILE, "r", encoding="utf-8") as f:
        template = json.load(f)

    master_doc = initialize_master_data(template)
    
    files = glob.glob(os.path.join(DOCS_DIR, "*"))
    print(f"Found {len(files)} files in {DOCS_DIR}")
    
    for filepath in files:
        print(f"Processing {filepath}...")
        content = extract_text_from_file(filepath)
        if not content.strip():
            print("  Skipping (empty or unreadable)")
            continue
            
        print(f"  Extracted {len(content)} chars.")
        extracted_data = process_file_with_ai(os.path.basename(filepath), content, template)
        
        # Update/Merge
        update_master_document(master_doc, extracted_data, template)
        
        # Update source reviewed metadata
        if "document_metadata" in master_doc and "sources_reviewed" in master_doc["document_metadata"]:
             current = master_doc["document_metadata"]["sources_reviewed"]
             fname = os.path.basename(filepath)
             if current:
                 if fname not in current:
                     master_doc["document_metadata"]["sources_reviewed"] += f", {fname}"
             else:
                 master_doc["document_metadata"]["sources_reviewed"] = fname

    # Write JSON
    with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
        json.dump(master_doc, f, indent=2)
    print(f"Saved JSON to {OUTPUT_JSON}")
    
    # Write DOCX
    generate_docx(master_doc, OUTPUT_DOCX)

if __name__ == "__main__":
    main()
